import json
import os
import time
import requests
import re
import sys
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# --------------------------------------------------------------------------
# 路径修复关键点：确保能导入位于项目根目录的 prompt 模块
# 如果你在 PyCharm 运行，通常它会自动处理；如果是命令行，这行代码能救命
current_script_path = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(current_script_path)  # 回退一级到 SawProject
sys.path.append(project_root)
# --------------------------------------------------------------------------

try:
    from prompt.pptx_prompt import build_pptx_layout_prompt
except ImportError:
    print("⚠️ 警告: 找不到 prompt 模块。请确保 prompt 文件夹在 SawProject 根目录下。")
    # 如果真的找不到，这里可以写一个简单的 fallback 或者直接退出
    # exit()

load_dotenv()


class PPTGenerator:
    def __init__(self, model="google/gemini-2.5-flash"):
        self.api_key = os.getenv("OPENROUTER_API_KEY")
        self.api_url = "https://openrouter.ai/api/v1/chat/completions"
        self.model = model
        self.output_file = os.path.join(project_root, "core", "Forest_Sports_Day.pptx")

        # 🟢 修复核心：assets 不在 core 里，而在上一级 (Project Root)
        # 逻辑：当前脚本(core/) -> 上一级(SawProject/) -> 拼接 assets
        self.assets_dir = os.path.join(project_root, "assets")

        if not self.api_key:
            # 尝试去根目录找 .env
            env_path = os.path.join(project_root, ".env")
            load_dotenv(env_path)
            self.api_key = os.getenv("OPENROUTER_API_KEY")

        if not self.api_key:
            raise ValueError("❌ 未找到 OPENROUTER_API_KEY，请检查 .env 文件")

        print(f"📂 资源目录设定为: {self.assets_dir}")
        self.file_index = self._index_assets()

    def _index_assets(self):
        """遍历 assets 目录，建立 {filename: full_path} 的映射"""
        index = {}
        if os.path.exists(self.assets_dir):
            for root, dirs, files in os.walk(self.assets_dir):
                for file in files:
                    index[file.lower()] = os.path.join(root, file)
        else:
            print(f"❌ 严重错误: 路径不存在 -> {self.assets_dir}")
        return index

    def find_real_image_path(self, llm_path):
        """智能查找路径"""
        # 1. 尝试直接路径拼接
        direct_path = os.path.join(self.assets_dir, llm_path)
        if os.path.exists(direct_path):
            return direct_path

        # 2. 尝试从文件名查找 (忽略 LLM 瞎编的目录结构)
        filename = os.path.basename(llm_path).lower()
        if filename in self.file_index:
            return self.file_index[filename]

        return None

    def load_data(self):
        # 同样，数据文件也要去 core 文件夹找（因为脚本在 core 里）
        # 或者如果它们也在根目录，要调整路径
        # 假设 json 文件和脚本都在 core 目录下：
        script_dir = os.path.dirname(os.path.abspath(__file__))

        narrative_path = os.path.join(script_dir, 'narrative_output.json')
        assets_list_path = os.path.join(script_dir, 'assets_list.json')

        try:
            with open(narrative_path, 'r', encoding='utf-8') as f:
                narrative = json.load(f)
            # 注意：assets_list.json 可能也在 core 下，或者在 assets 下？
            # 根据你的截图，assets_list.json 在 core 文件夹里，所以上面路径是对的
            with open(assets_list_path, 'r', encoding='utf-8') as f:
                assets = json.load(f)
            return narrative, assets
        except FileNotFoundError as e:
            print(f"❌ 数据文件丢失: {e}")
            return None, None

    def call_llm_layout(self, narrative_chunk, assets_full):
        # 检查 build_pptx_layout_prompt 是否成功导入
        if 'build_pptx_layout_prompt' not in globals():
            print("❌ 无法生成 Prompt，因为导入失败")
            return None

        prompt = build_pptx_layout_prompt(narrative_chunk, assets_full)

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
            "HTTP-Referer": "https://github.com/YourProject",
            "X-Title": "PPTX Generator"
        }

        payload = {
            "model": self.model,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.2,
            "response_format": {"type": "json_object"}
        }

        print(f"📡 请求 LLM 生成 {len(narrative_chunk)} 页的布局数据...")
        try:
            response = requests.post(self.api_url, headers=headers, json=payload, timeout=60)
            response.raise_for_status()
            content = response.json()['choices'][0]['message']['content']
            return self.clean_json_response(content)
        except Exception as e:
            print(f"❌ API 请求失败: {e}")
            return None

    def clean_json_response(self, text):
        text = re.sub(r'^```json\s*', '', text, flags=re.MULTILINE)
        text = re.sub(r'^```\s*', '', text, flags=re.MULTILINE)
        text = text.strip()
        try:
            return json.loads(text)
        except json.JSONDecodeError as e:
            print(f"❌ JSON 解析失败: {e}")
            return None

    def create_slide(self, prs, slide_data):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        page_id = slide_data.get('page_id', 'Unknown')
        print(f"  🎨 正在绘制第 {page_id} 页...")

        elements = slide_data.get('elements', [])

        for el in elements:
            el_type = el.get('type')
            pos = el.get('position', {})

            left = Inches(pos.get('left', 0))
            top = Inches(pos.get('top', 0))
            width = Inches(pos.get('width', 1))
            height = Inches(pos.get('height', 1))

            if el_type == 'image':
                raw_path = el.get('content', '')
                real_path = self.find_real_image_path(raw_path)

                if real_path:
                    try:
                        slide.shapes.add_picture(real_path, left, top, width, height)
                    except Exception as e:
                        print(f"    ⚠️ 图片加载出错: {os.path.basename(real_path)} ({e})")
                else:
                    print(f"    ⚠️ 警告: 图片未找到 (LLM请求: {raw_path})")

            elif el_type == 'shape':
                shape_type = MSO_SHAPE.RECTANGLE
                if el.get('shape_type') == 'ROUNDED_RECTANGLE':
                    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE

                shape = slide.shapes.add_shape(shape_type, left, top, width, height)

                fill = shape.fill
                fill.solid()
                hex_color = el.get('color_hex', 'FFFFFF')
                try:
                    fill.fore_color.rgb = RGBColor.from_string(hex_color)
                except:
                    fill.fore_color.rgb = RGBColor(255, 255, 255)

                shape.line.fill.background()

            elif el_type == 'text':
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.text = el.get('content', '')
                p.font.size = Pt(el.get('font_size', 18))

                font_color = el.get('font_color', '000000')
                try:
                    p.font.color.rgb = RGBColor.from_string(font_color)
                except:
                    pass
                p.alignment = PP_ALIGN.LEFT

    def generate(self):
        narrative, assets = self.load_data()
        if not narrative: return

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        chunk_size = 4
        chunks = [narrative[i:i + chunk_size] for i in range(0, len(narrative), chunk_size)]

        for index, chunk in enumerate(chunks):
            retry_count = 0
            max_retries = 2
            success = False

            while retry_count < max_retries and not success:
                try:
                    layout_data_list = self.call_llm_layout(chunk, assets)
                    if layout_data_list:
                        if isinstance(layout_data_list, list):
                            for slide_data in layout_data_list:
                                self.create_slide(prs, slide_data)
                        else:
                            self.create_slide(prs, layout_data_list)
                        success = True
                    else:
                        retry_count += 1
                        print(f"    ⚠️ LLM 返回数据为空，重试 {retry_count}/{max_retries}...")
                except Exception as e:
                    print(f"    ❌ 批次处理发生异常: {e}")
                    retry_count += 1

            time.sleep(2)

        prs.save(self.output_file)
        print(f"🎉 PPT 生成完成: {self.output_file}")


if __name__ == "__main__":
    generator = PPTGenerator()
    generator.generate()